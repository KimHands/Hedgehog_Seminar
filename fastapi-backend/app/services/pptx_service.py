import io
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont


def extract_slides(pptx_bytes: bytes) -> list[dict]:
    """PPTX에서 슬라이드별 텍스트를 추출합니다."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line:
                        texts.append(line)
        slides.append({
            "slide_number": idx,
            "text_content": "\n".join(texts),
        })

    return slides


def generate_thumbnail(pptx_bytes: bytes, slide_number: int, width: int = 960, height: int = 540) -> bytes:
    """슬라이드 썸네일 이미지를 생성합니다."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide = prs.slides[slide_number - 1]

    img = Image.new("RGB", (width, height), color=(30, 58, 95))
    draw = ImageDraw.Draw(img)

    # 슬라이드 텍스트를 이미지에 렌더링 (간단한 구현)
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in para.runs if run.text.strip())
                if line:
                    texts.append(line)

    y = 40
    for text in texts[:8]:
        draw.text((40, y), text[:60], fill=(255, 255, 255))
        y += 36

    draw.text((width - 80, height - 30), f"{slide_number}", fill=(200, 200, 200))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()
